#!/usr/bin/env python3
"""
One-command slide export: PPTX to PNG images.

Auto-detects the best available method for your platform:

    macOS:   Keynote AppleScript (pixel-perfect, pre-installed on all Macs)
    Windows: PowerPoint COM (perfect, via pywin32)
    All:     python-pptx + Pillow (fast, basic quality, always available)

Usage:
    python3 export_slides.py presentation.pptx slides/
"""

import os, sys, re, time, subprocess, platform, shutil, argparse
from pathlib import Path


def has_keynote():
    return platform.system() == "Darwin" and os.path.exists("/Applications/Keynote.app")


def has_powerpoint_win():
    if platform.system() != "Windows":
        return False
    try:
        import win32com.client
        return True
    except ImportError:
        return False


def export_with_keynote(pptx_path, output_dir):
    """Export PPTX slides as PNG using Keynote via AppleScript."""
    abs_pptx = os.path.abspath(pptx_path)
    abs_out = os.path.abspath(output_dir)
    os.makedirs(abs_out, exist_ok=True)
    
    subprocess.run(["killall", "Keynote"], capture_output=True)
    time.sleep(1)
    
    tmp_prefix = os.path.join(abs_out, "_tmp_slide")
    
    script = (
        'tell application "Keynote"\n'
        '    launch\n'
        '    delay 3\n'
        f'    set theDoc to open POSIX file "{abs_pptx}"\n'
        '    delay 3\n'
        '    try\n'
        f'        export theDoc to POSIX file "{tmp_prefix}" as slide images with properties {{image format:PNG, compression factor:1.0}}\n'
        '        return "SUCCESS"\n'
        '    on error errMsg\n'
        '        return "KEYNOTE_ERROR: " & errMsg\n'
        '    end try\n'
        'end tell'
    )
    
    print("  Exporting via Keynote (~10 seconds)...")
    result = subprocess.run(["osascript", "-e", script],
                          capture_output=True, text=True, timeout=120)
    
    output = result.stdout.strip()
    print(f"  Keynote: {output[:100]}")
    
    if "KEYNOTE_ERROR" in output:
        print(f"  Keynote export failed: {output}")
        subprocess.run(["osascript", "-e", 'tell application "Keynote" to quit'],
                      capture_output=True)
        return False
    
    # Rename: _tmp_slide.001.png -> Slide01.png
    tmp_path = Path(abs_out)
    png_files = list(tmp_path.glob("_tmp_slide.*.png"))
    
    subfolder = tmp_path / "_tmp_slide"
    if subfolder.is_dir():
        png_files.extend(subfolder.glob("*.png"))
    
    count = 0
    for f in sorted(png_files):
        match = re.search(r'[._](\d+)\.png$', f.name)
        if match:
            num = int(match.group(1))
            new_name = tmp_path / f"Slide{num:02d}.png"
            f.rename(new_name)
            size_kb = new_name.stat().st_size / 1024
            print(f"  Slide {num:02d}: {size_kb:.0f} KB")
            count += 1
    
    # Cleanup
    for d in tmp_path.glob("_tmp_slide*"):
        if d.is_dir():
            shutil.rmtree(d)
        elif d.is_file():
            d.unlink()
    
    subprocess.run(["osascript", "-e", 'tell application "Keynote" to quit'],
                  capture_output=True)
    
    print(f"\n  {count} slides exported via Keynote")
    return count > 0


def export_with_powerpoint_win(pptx_path, output_dir):
    """Export PPTX slides as PNG using PowerPoint COM on Windows."""
    import win32com.client
    
    abs_pptx = os.path.abspath(pptx_path)
    abs_out = os.path.abspath(output_dir)
    os.makedirs(abs_out, exist_ok=True)
    
    powerpoint = win32com.client.Dispatch("PowerPoint.Application")
    presentation = powerpoint.Presentations.Open(abs_pptx, WithWindow=False)
    
    slide_w = presentation.PageSetup.SlideWidth
    slide_h = presentation.PageSetup.SlideHeight
    scale = 2
    
    count = 0
    for i in range(1, presentation.Slides.Count + 1):
        slide = presentation.Slides(i)
        image_path = os.path.join(abs_out, f"Slide{i:02d}.png")
        slide.Export(image_path, "PNG", int(slide_w * scale), int(slide_h * scale))
        count += 1
        print(f"  Slide {i:02d}: exported")
    
    presentation.Close()
    print(f"\n  {count} slides exported via PowerPoint")
    return count > 0


def export_with_python(pptx_path, output_dir, dpi=150):
    """Fallback: export using python-pptx + Pillow."""
    from pptx import Presentation
    from PIL import Image, ImageDraw, ImageFont
    import io
    
    prs = Presentation(pptx_path)
    os.makedirs(output_dir, exist_ok=True)
    
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    px_w = int(slide_w * dpi / 914400)
    px_h = int(slide_h * dpi / 914400)
    
    total = len(list(prs.slides))
    
    font_paths = [
        "/System/Library/Fonts/Helvetica.ttc",
        "/System/Library/Fonts/SFNSText.ttf",
        "/System/Library/Fonts/Supplemental/Arial.ttf",
        "/Library/Fonts/Arial.ttf",
        "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
    ]
    font_path = next((fp for fp in font_paths if os.path.exists(fp)), None)
    
    for i, slide in enumerate(prs.slides):
        img = Image.new("RGB", (px_w, px_h), (255, 255, 255))
        draw = ImageDraw.Draw(img)
        
        for shape in slide.shapes:
            x = int(shape.left * dpi / 914400)
            y = int(shape.top * dpi / 914400)
            w = int(shape.width * dpi / 914400)
            h = int(shape.height * dpi / 914400)
            
            if shape.shape_type == 13:
                try:
                    pil_img = Image.open(io.BytesIO(shape.image.blob))
                    pil_img = pil_img.resize((w, h), Image.LANCZOS)
                    if pil_img.mode == "RGBA":
                        img.paste(pil_img, (x, y), pil_img)
                    else:
                        img.paste(pil_img, (x, y))
                except:
                    draw.rectangle([x, y, x+w, y+h], outline=(200, 200, 200))
                continue
            
            if shape.has_text_frame:
                default_size = 12
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            pts = run.font.size / 12700
                            if pts > 4:
                                default_size = int(pts)
                            break
                try:
                    font = ImageFont.truetype(font_path, max(default_size, 9)) if font_path else ImageFont.load_default()
                except:
                    font = ImageFont.load_default()
                
                draw.rectangle([x, y, x+w, y+h], fill=(250, 251, 252), outline=(220, 225, 230))
                line_y = y + 4
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if not text:
                        continue
                    chars_per_line = max(int(w / (default_size * 0.6)), 20)
                    words = text.split()
                    cur = ""
                    for word in words:
                        test = cur + (" " if cur else "") + word
                        if len(test) > chars_per_line:
                            draw.text((x+6, line_y), cur, fill=(30, 30, 30), font=font)
                            line_y += default_size + 4
                            cur = word
                        else:
                            cur = test
                    if cur:
                        draw.text((x+6, line_y), cur, fill=(30, 30, 30), font=font)
                        line_y += default_size + 4
        
        out_path = os.path.join(output_dir, f"Slide{i+1:02d}.png")
        img.save(out_path, "PNG", optimize=True)
        print(f"  Slide {i+1:02d}/{total:02d}: {px_w}x{px_h} ({os.path.getsize(out_path)/1024:.0f} KB)")
    
    print(f"\n  {total} slides exported via python-pptx")


def main():
    parser = argparse.ArgumentParser(
        description="Convert PPTX slides to PNG images — auto-selects best method",
    )
    parser.add_argument("pptx", help="Path to PPTX file")
    parser.add_argument("output", help="Output directory for PNG images")
    parser.add_argument("--dpi", type=int, default=150, help="DPI for python-pptx (default: 150)")
    parser.add_argument("--force-python", action="store_true", help="Force python-pptx renderer")
    
    args = parser.parse_args()
    
    if not os.path.exists(args.pptx):
        print(f"File not found: {args.pptx}")
        sys.exit(1)
    
    print("=" * 50)
    print("Slide Exporter")
    print("=" * 50)
    print(f"\nInput:  {args.pptx}")
    print(f"Output: {args.output}/")
    
    if args.force_python:
        print("Method: python-pptx (forced)\n")
        export_with_python(args.pptx, args.output, args.dpi)
    elif has_keynote():
        print("Method: Keynote (pixel-perfect)\n")
        ok = export_with_keynote(args.pptx, args.output)
        if not ok:
            print("\nKeynote failed. Falling back to python-pptx...\n")
            export_with_python(args.pptx, args.output, args.dpi)
    elif has_powerpoint_win():
        print("Method: PowerPoint COM (perfect)\n")
        try:
            export_with_powerpoint_win(args.pptx, args.output)
        except Exception as e:
            print(f"\nPowerPoint failed: {e}")
            print("Falling back to python-pptx...\n")
            export_with_python(args.pptx, args.output, args.dpi)
    else:
        print("Method: python-pptx + Pillow (basic quality)\n")
        export_with_python(args.pptx, args.output, args.dpi)
    
    count = len(list(Path(args.output).glob("Slide*.png")))
    print(f"\nDone: {count} slides in {args.output}/")


if __name__ == "__main__":
    main()
