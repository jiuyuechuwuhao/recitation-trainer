#!/usr/bin/env python3
"""
Extract text content from a PPTX file to assist AI in generating director scripts.

Usage:
    python3 extract_pptx_text.py presentation.pptx

Output: Prints each slide's text content to stdout, grouped by slide.
The AI agent reads this output and writes the director script based on it.
"""

import sys
from pptx import Presentation


def extract_text(pptx_path):
    prs = Presentation(pptx_path)
    for i, slide in enumerate(prs.slides):
        print(f"\n{'='*60}")
        print(f"## Slide {i+1} — (AI: fill in title based on content below)")
        print(f"{'='*60}")
        
        # Collect all text, grouping by shape for context
        for j, shape in enumerate(slide.shapes):
            if shape.has_text_frame:
                texts = []
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
                if texts:
                    shape_label = f"[Shape {j+1}]"
                    for t in texts:
                        print(f"  {shape_label} {t}")
                        shape_label = "         "  # indent continuation lines
        
        # Check for images, videos, tables
        img_count = 0
        video_count = 0
        table_count = 0
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                img_count += 1
            elif shape.shape_type == 15:  # Media
                video_count += 1
            elif shape.has_table:
                table_count += 1
        
        extras = []
        if img_count: extras.append(f"{img_count} image(s)")
        if video_count: extras.append(f"{video_count} video(s)")
        if table_count: extras.append(f"{table_count} table(s)")
        if extras:
            print(f"  [Contains: {', '.join(extras)}]")

    print(f"\n{'='*60}")
    print(f"Total slides: {len(prs.slides)}")
    print(f"{'='*60}")


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 extract_pptx_text.py <presentation.pptx>")
        sys.exit(1)
    extract_text(sys.argv[1])
